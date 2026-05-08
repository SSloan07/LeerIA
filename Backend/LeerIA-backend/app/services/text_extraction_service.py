from pathlib import Path
import tempfile

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx"}


def get_file_extension(file_name: str) -> str:
    return Path(file_name).suffix.lower()


def validate_supported_file(file_name: str) -> None:
    extension = get_file_extension(file_name)

    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Tipo de archivo no soportado: {extension}")


def extract_text_from_txt(file_path: str) -> str:
    path = Path(file_path)

    return path.read_text(encoding="utf-8", errors="ignore")


def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)

    pages_text: list[str] = []

    for page_number, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""

        if page_text.strip():
            pages_text.append(f"\n--- Página {page_number} ---\n{page_text}")

    return "\n".join(pages_text).strip()


def extract_text_from_docx(file_path: str) -> str:
    document = Document(file_path)

    paragraphs = [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]

    return "\n".join(paragraphs).strip()


def extract_text_from_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)

    slides_text: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_parts.append(shape.text.strip())

        if slide_parts:
            slides_text.append(
                f"\n--- Diapositiva {slide_index} ---\n" + "\n".join(slide_parts)
            )

    return "\n".join(slides_text).strip()


def extract_text_from_file(file_path: str, file_name: str | None = None) -> str:
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"No existe el archivo: {file_path}")

    name_to_validate = file_name or path.name
    validate_supported_file(name_to_validate)

    extension = get_file_extension(name_to_validate)

    if extension == ".txt":
        return extract_text_from_txt(file_path)

    if extension == ".pdf":
        return extract_text_from_pdf(file_path)

    if extension == ".docx":
        return extract_text_from_docx(file_path)

    if extension == ".pptx":
        return extract_text_from_pptx(file_path)

    raise ValueError(f"No hay extractor disponible para: {extension}")


def extract_text_from_bytes(file_bytes: bytes, file_name: str) -> str:
    validate_supported_file(file_name)

    extension = get_file_extension(file_name)

    with tempfile.NamedTemporaryFile(delete=True, suffix=extension) as temp_file:
        temp_file.write(file_bytes)
        temp_file.flush()

        return extract_text_from_file(
            file_path=temp_file.name,
            file_name=file_name,
        )